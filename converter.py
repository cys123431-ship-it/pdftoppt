import csv
import io
import os
import tempfile
from datetime import datetime
from typing import Callable, Optional, Sequence

import fitz  # PyMuPDF
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches

ProgressCallback = Optional[Callable[[int], None]]

CONFLICT_OVERWRITE = "overwrite"
CONFLICT_SKIP = "skip"
CONFLICT_AUTO_RENAME = "auto_rename"
VALID_CONFLICT_POLICIES = {
    CONFLICT_OVERWRITE,
    CONFLICT_SKIP,
    CONFLICT_AUTO_RENAME,
}

CANCELLED_MESSAGE = "Cancelled by user."


def _set_progress(progress_callback: ProgressCallback, value: float) -> None:
    if progress_callback:
        clamped = max(0, min(100, int(value)))
        progress_callback(clamped)


def _is_cancelled(cancel_event: object | None) -> bool:
    return bool(cancel_event and hasattr(cancel_event, "is_set") and cancel_event.is_set())


def _validate_conflict_policy(output_conflict_policy: str) -> None:
    if output_conflict_policy not in VALID_CONFLICT_POLICIES:
        raise ValueError(f"Unsupported conflict policy: {output_conflict_policy}")


def _resolve_output_path(output_path: str, output_conflict_policy: str) -> tuple[str, bool]:
    _validate_conflict_policy(output_conflict_policy)
    if not os.path.exists(output_path):
        return output_path, False

    if output_conflict_policy == CONFLICT_OVERWRITE:
        return output_path, False

    if output_conflict_policy == CONFLICT_SKIP:
        return "", True

    directory, filename = os.path.split(output_path)
    stem, extension = os.path.splitext(filename)
    suffix = 1
    while True:
        candidate = os.path.join(directory, f"{stem}_{suffix}{extension}")
        if not os.path.exists(candidate):
            return candidate, False
        suffix += 1


def _resolve_output_directory(output_dir: str, output_conflict_policy: str) -> tuple[str, bool]:
    _validate_conflict_policy(output_conflict_policy)
    if not os.path.exists(output_dir):
        return output_dir, False

    if not os.path.isdir(output_dir):
        raise ValueError(f"Output directory path is an existing file: {output_dir}")

    if output_conflict_policy == CONFLICT_OVERWRITE:
        return output_dir, False

    if output_conflict_policy == CONFLICT_SKIP:
        return "", True

    suffix = 1
    while True:
        candidate = f"{output_dir}_{suffix}"
        if not os.path.exists(candidate):
            return candidate, False
        suffix += 1


def _remove_existing_file(path: str) -> None:
    if os.path.exists(path) and os.path.isfile(path):
        os.remove(path)


def _build_output_note(original_path: str, resolved_path: str) -> str:
    if os.path.normcase(os.path.abspath(original_path)) == os.path.normcase(os.path.abspath(resolved_path)):
        return ""
    return f" Saved as: {resolved_path}"


def parse_page_range(page_range_text: str, total_pages: int) -> list[int]:
    """
    Parses page ranges like "1-3,5,8-10" into zero-based page indices.
    Empty input means all pages.
    """
    if total_pages <= 0:
        return []

    if not page_range_text or not page_range_text.strip():
        return list(range(total_pages))

    selected_pages: set[int] = set()
    tokens = page_range_text.split(",")

    for token in tokens:
        token = token.strip()
        if not token:
            continue

        if "-" in token:
            parts = [part.strip() for part in token.split("-", 1)]
            if len(parts) != 2 or not parts[0].isdigit() or not parts[1].isdigit():
                raise ValueError(f"Invalid page range token: {token}")

            start_page = int(parts[0])
            end_page = int(parts[1])
            if start_page > end_page:
                raise ValueError(f"Range start must be <= end: {token}")
            if start_page < 1 or end_page > total_pages:
                raise ValueError(f"Page range out of bounds: {token} (valid: 1-{total_pages})")

            selected_pages.update(range(start_page - 1, end_page))
            continue

        if not token.isdigit():
            raise ValueError(f"Invalid page number: {token}")

        page_number = int(token)
        if page_number < 1 or page_number > total_pages:
            raise ValueError(f"Page out of bounds: {page_number} (valid: 1-{total_pages})")
        selected_pages.add(page_number - 1)

    if not selected_pages:
        raise ValueError("No pages selected.")

    return sorted(selected_pages)


def _open_pdf_document(pdf_path: str, input_password: str = "") -> fitz.Document:
    document = fitz.open(pdf_path)
    if document.needs_pass:
        if not input_password:
            document.close()
            raise ValueError("This PDF is password-protected. Enter an input password.")

        authenticated = document.authenticate(input_password)
        if authenticated == 0:
            document.close()
            raise ValueError("Invalid password for encrypted PDF.")

    return document


def _save_pdf_document(document: fitz.Document, output_pdf_path: str, output_password: str = "") -> None:
    if output_password:
        encryption_mode = getattr(fitz, "PDF_ENCRYPT_AES_256", fitz.PDF_ENCRYPT_AES_128)
        document.save(
            output_pdf_path,
            encryption=encryption_mode,
            owner_pw=output_password,
            user_pw=output_password,
        )
    else:
        document.save(output_pdf_path)


def _create_temp_pdf_with_selected_pages(
    source_doc: fitz.Document,
    page_indices: Sequence[int],
    cancel_event: object | None = None,
) -> str:
    temp_doc = fitz.open()
    try:
        for page_index in page_indices:
            if _is_cancelled(cancel_event):
                raise RuntimeError(CANCELLED_MESSAGE)
            temp_doc.insert_pdf(source_doc, from_page=page_index, to_page=page_index)

        file_descriptor, temp_pdf_path = tempfile.mkstemp(suffix=".pdf")
        os.close(file_descriptor)
        temp_doc.save(temp_pdf_path)
        return temp_pdf_path
    finally:
        temp_doc.close()


def _write_batch_failure_log(output_dir: str, rows: Sequence[tuple[str, str, str]]) -> str:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    log_path = os.path.join(output_dir, f"batch_failures_{timestamp}.csv")
    with open(log_path, "w", newline="", encoding="utf-8-sig") as output_file:
        writer = csv.writer(output_file)
        writer.writerow(["file", "status", "message"])
        for file_name, status, message in rows:
            writer.writerow([file_name, status, message])
    return log_path


def convert_pdf_to_pptx(
    pdf_path: str,
    pptx_path: str,
    progress_callback: ProgressCallback = None,
    page_range_text: str = "",
    input_password: str = "",
    output_conflict_policy: str = CONFLICT_OVERWRITE,
    render_dpi: int = 144,
    cancel_event: object | None = None,
) -> tuple[bool, str]:
    """Converts selected pages from PDF to PPTX (one page image per slide)."""
    doc = None
    try:
        if render_dpi <= 0:
            return False, "Render DPI must be greater than 0."
        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        resolved_pptx_path, skipped = _resolve_output_path(pptx_path, output_conflict_policy)
        if skipped:
            return True, f"Skipped existing file: {pptx_path}"
        if output_conflict_policy == CONFLICT_OVERWRITE:
            _remove_existing_file(resolved_pptx_path)

        doc = _open_pdf_document(pdf_path, input_password)
        selected_pages = parse_page_range(page_range_text, len(doc))
        if not selected_pages:
            return False, "The PDF has no pages."

        prs = Presentation()
        first_page = doc[selected_pages[0]]
        prs.slide_width = Inches(first_page.rect.width / 72.0)
        prs.slide_height = Inches(first_page.rect.height / 72.0)

        zoom = render_dpi / 72.0
        total_pages = len(selected_pages)
        for index, page_number in enumerate(selected_pages):
            if _is_cancelled(cancel_event):
                return False, CANCELLED_MESSAGE

            page = doc[page_number]
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            img_data = pix.tobytes("png")
            image_stream = io.BytesIO(img_data)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                image_stream,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            _set_progress(progress_callback, ((index + 1) / total_pages) * 100)

        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        prs.save(resolved_pptx_path)
        output_note = _build_output_note(pptx_path, resolved_pptx_path)
        return True, f"Conversion successful!{output_note}"
    except Exception as exc:
        return False, str(exc)
    finally:
        if doc is not None:
            doc.close()


def convert_pdf_to_docx(
    pdf_path: str,
    docx_path: str,
    progress_callback: ProgressCallback = None,
    page_range_text: str = "",
    input_password: str = "",
    output_conflict_policy: str = CONFLICT_OVERWRITE,
    cancel_event: object | None = None,
) -> tuple[bool, str]:
    """Converts selected pages from PDF to DOCX."""
    source_doc = None
    cv = None
    temp_pdf_path = None

    try:
        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        resolved_docx_path, skipped = _resolve_output_path(docx_path, output_conflict_policy)
        if skipped:
            return True, f"Skipped existing file: {docx_path}"
        if output_conflict_policy == CONFLICT_OVERWRITE:
            _remove_existing_file(resolved_docx_path)

        source_doc = _open_pdf_document(pdf_path, input_password)
        selected_pages = parse_page_range(page_range_text, len(source_doc))
        if not selected_pages:
            return False, "The PDF has no pages."
        _set_progress(progress_callback, 10)

        temp_pdf_path = _create_temp_pdf_with_selected_pages(source_doc, selected_pages, cancel_event)
        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        _set_progress(progress_callback, 35)
        cv = Converter(temp_pdf_path)
        cv.convert(resolved_docx_path)
        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        _set_progress(progress_callback, 100)
        output_note = _build_output_note(docx_path, resolved_docx_path)
        return True, f"Conversion successful!{output_note}"
    except RuntimeError as exc:
        if str(exc) == CANCELLED_MESSAGE:
            return False, CANCELLED_MESSAGE
        return False, str(exc)
    except Exception as exc:
        return False, str(exc)
    finally:
        if cv is not None:
            cv.close()
        if source_doc is not None:
            source_doc.close()
        if temp_pdf_path and os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)


def convert_pdf_to_images(
    pdf_path: str,
    output_dir: str,
    image_format: str = "png",
    dpi: int = 144,
    progress_callback: ProgressCallback = None,
    page_range_text: str = "",
    input_password: str = "",
    output_conflict_policy: str = CONFLICT_OVERWRITE,
    jpg_quality: int = 90,
    cancel_event: object | None = None,
) -> tuple[bool, str]:
    """Converts selected PDF pages to PNG or JPG files."""
    doc = None
    try:
        if dpi <= 0:
            return False, "DPI must be greater than 0."
        if jpg_quality < 1 or jpg_quality > 100:
            return False, "JPG quality must be between 1 and 100."
        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        normalized_format = image_format.lower()
        if normalized_format not in {"png", "jpg", "jpeg"}:
            return False, f"Unsupported image format: {image_format}"

        output_extension = "jpg" if normalized_format in {"jpg", "jpeg"} else "png"
        os.makedirs(output_dir, exist_ok=True)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]

        doc = _open_pdf_document(pdf_path, input_password)
        selected_pages = parse_page_range(page_range_text, len(doc))
        if not selected_pages:
            return False, "The PDF has no pages."

        zoom = dpi / 72.0
        total_pages = len(selected_pages)
        created_count = 0
        skipped_count = 0

        for index, page_number in enumerate(selected_pages):
            if _is_cancelled(cancel_event):
                return False, CANCELLED_MESSAGE

            output_name = f"{base_name}_p{page_number + 1:03d}.{output_extension}"
            output_path = os.path.join(output_dir, output_name)
            resolved_path, skipped = _resolve_output_path(output_path, output_conflict_policy)

            if skipped:
                skipped_count += 1
                _set_progress(progress_callback, ((index + 1) / total_pages) * 100)
                continue

            if output_conflict_policy == CONFLICT_OVERWRITE:
                _remove_existing_file(resolved_path)

            page = doc[page_number]
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            if output_extension == "jpg":
                image_bytes = pix.tobytes("jpeg", jpg_quality=jpg_quality)
            else:
                image_bytes = pix.tobytes("png")

            with open(resolved_path, "wb") as output_file:
                output_file.write(image_bytes)

            created_count += 1
            _set_progress(progress_callback, ((index + 1) / total_pages) * 100)

        if created_count == 0 and skipped_count > 0:
            return True, "All images were skipped because output files already exist."
        if skipped_count > 0:
            return True, f"Saved {created_count} images. Skipped {skipped_count} existing files."
        return True, f"Saved {created_count} image files."
    except Exception as exc:
        return False, str(exc)
    finally:
        if doc is not None:
            doc.close()


def merge_pdfs(
    input_pdf_paths: Sequence[str],
    output_pdf_path: str,
    progress_callback: ProgressCallback = None,
    input_password: str = "",
    output_password: str = "",
    output_conflict_policy: str = CONFLICT_OVERWRITE,
    cancel_event: object | None = None,
) -> tuple[bool, str]:
    """Merges multiple PDFs into one file."""
    if len(input_pdf_paths) < 2:
        return False, "Select at least 2 PDF files to merge."

    merged_doc = fitz.open()
    try:
        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        resolved_output_path, skipped = _resolve_output_path(output_pdf_path, output_conflict_policy)
        if skipped:
            return True, f"Skipped existing file: {output_pdf_path}"
        if output_conflict_policy == CONFLICT_OVERWRITE:
            _remove_existing_file(resolved_output_path)

        total_files = len(input_pdf_paths)
        for index, input_path in enumerate(input_pdf_paths):
            if _is_cancelled(cancel_event):
                return False, CANCELLED_MESSAGE

            source_doc = _open_pdf_document(input_path, input_password)
            try:
                merged_doc.insert_pdf(source_doc)
            finally:
                source_doc.close()

            _set_progress(progress_callback, ((index + 1) / total_files) * 100)

        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        _save_pdf_document(merged_doc, resolved_output_path, output_password)
        output_note = _build_output_note(output_pdf_path, resolved_output_path)
        return True, f"Merge successful!{output_note}"
    except Exception as exc:
        return False, str(exc)
    finally:
        merged_doc.close()


def split_pdf(
    pdf_path: str,
    output_dir: str,
    progress_callback: ProgressCallback = None,
    page_range_text: str = "",
    input_password: str = "",
    output_password: str = "",
    output_conflict_policy: str = CONFLICT_OVERWRITE,
    cancel_event: object | None = None,
) -> tuple[bool, str]:
    """Splits selected pages into one-page PDF files."""
    source_doc = None
    try:
        if _is_cancelled(cancel_event):
            return False, CANCELLED_MESSAGE

        os.makedirs(output_dir, exist_ok=True)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]

        source_doc = _open_pdf_document(pdf_path, input_password)
        selected_pages = parse_page_range(page_range_text, len(source_doc))
        if not selected_pages:
            return False, "The PDF has no pages."

        total_pages = len(selected_pages)
        created_count = 0
        skipped_count = 0

        for index, page_number in enumerate(selected_pages):
            if _is_cancelled(cancel_event):
                return False, CANCELLED_MESSAGE

            output_name = f"{base_name}_p{page_number + 1:03d}.pdf"
            output_path = os.path.join(output_dir, output_name)
            resolved_output_path, skipped = _resolve_output_path(output_path, output_conflict_policy)
            if skipped:
                skipped_count += 1
                _set_progress(progress_callback, ((index + 1) / total_pages) * 100)
                continue

            if output_conflict_policy == CONFLICT_OVERWRITE:
                _remove_existing_file(resolved_output_path)

            split_doc = fitz.open()
            try:
                split_doc.insert_pdf(source_doc, from_page=page_number, to_page=page_number)
                _save_pdf_document(split_doc, resolved_output_path, output_password)
            finally:
                split_doc.close()

            created_count += 1
            _set_progress(progress_callback, ((index + 1) / total_pages) * 100)

        if created_count == 0 and skipped_count > 0:
            return True, "All split files were skipped because output files already exist."
        if skipped_count > 0:
            return True, f"Created {created_count} split files. Skipped {skipped_count} existing files."
        return True, f"Created {created_count} split PDF files."
    except Exception as exc:
        return False, str(exc)
    finally:
        if source_doc is not None:
            source_doc.close()


def batch_convert_folder(
    input_dir: str,
    output_dir: str,
    target_format: str,
    progress_callback: ProgressCallback = None,
    page_range_text: str = "",
    input_password: str = "",
    output_password: str = "",
    output_conflict_policy: str = CONFLICT_OVERWRITE,
    render_dpi: int = 144,
    jpg_quality: int = 90,
    write_failure_log: bool = True,
    cancel_event: object | None = None,
) -> tuple[bool, str]:
    """
    Batch converts all PDFs in a folder into target format.
    Supported target formats: PPTX, DOCX, PNG, JPG
    """
    try:
        if not os.path.isdir(input_dir):
            return False, "Input folder does not exist."

        os.makedirs(output_dir, exist_ok=True)
        pdf_files = sorted(
            filename
            for filename in os.listdir(input_dir)
            if filename.lower().endswith(".pdf") and os.path.isfile(os.path.join(input_dir, filename))
        )

        if not pdf_files:
            return False, "No PDF files found in input folder."

        normalized_target = target_format.upper()
        if normalized_target not in {"PPTX", "DOCX", "PNG", "JPG"}:
            return False, f"Unsupported target format: {target_format}"

        total_files = len(pdf_files)
        converted_count = 0
        skipped_count = 0
        failed_count = 0
        failure_rows: list[tuple[str, str, str]] = []
        cancelled = False

        for file_index, filename in enumerate(pdf_files):
            if _is_cancelled(cancel_event):
                cancelled = True
                failure_rows.append((filename, "cancelled", CANCELLED_MESSAGE))
                break

            input_pdf_path = os.path.join(input_dir, filename)
            base_name = os.path.splitext(filename)[0]

            def file_progress(child_percent: int, index: int = file_index) -> None:
                overall = ((index + (child_percent / 100.0)) / total_files) * 100.0
                _set_progress(progress_callback, overall)

            if normalized_target == "PPTX":
                output_path = os.path.join(output_dir, f"{base_name}.pptx")
                success, message = convert_pdf_to_pptx(
                    input_pdf_path,
                    output_path,
                    file_progress,
                    page_range_text,
                    input_password,
                    output_conflict_policy,
                    render_dpi,
                    cancel_event,
                )
            elif normalized_target == "DOCX":
                output_path = os.path.join(output_dir, f"{base_name}.docx")
                success, message = convert_pdf_to_docx(
                    input_pdf_path,
                    output_path,
                    file_progress,
                    page_range_text,
                    input_password,
                    output_conflict_policy,
                    cancel_event,
                )
            else:
                output_subdir = os.path.join(output_dir, base_name)
                resolved_subdir, dir_skipped = _resolve_output_directory(
                    output_subdir,
                    output_conflict_policy,
                )
                if dir_skipped:
                    success, message = True, f"Skipped existing output directory: {output_subdir}"
                else:
                    os.makedirs(resolved_subdir, exist_ok=True)
                    success, message = convert_pdf_to_images(
                        input_pdf_path,
                        resolved_subdir,
                        normalized_target.lower(),
                        render_dpi,
                        file_progress,
                        page_range_text,
                        input_password,
                        output_conflict_policy,
                        jpg_quality,
                        cancel_event,
                    )

            if not success and message == CANCELLED_MESSAGE:
                cancelled = True
                failure_rows.append((filename, "cancelled", message))
                break

            if success and message.startswith("Skipped"):
                skipped_count += 1
                failure_rows.append((filename, "skipped", message))
            elif success:
                converted_count += 1
            else:
                failed_count += 1
                failure_rows.append((filename, "failed", message))

            _set_progress(progress_callback, ((file_index + 1) / total_files) * 100)

        log_suffix = ""
        if write_failure_log and failure_rows:
            log_path = _write_batch_failure_log(output_dir, failure_rows)
            log_suffix = f" Failure log: {log_path}"

        if cancelled:
            return (
                False,
                f"{CANCELLED_MESSAGE} Converted {converted_count}/{total_files} files. "
                f"Failed: {failed_count}, Skipped: {skipped_count}.{log_suffix}",
            )

        if failed_count > 0:
            return (
                False,
                f"Completed with errors. Converted {converted_count}/{total_files}, "
                f"Failed: {failed_count}, Skipped: {skipped_count}.{log_suffix}",
            )

        if skipped_count > 0:
            return (
                True,
                f"Completed with skips. Converted {converted_count}/{total_files}, "
                f"Skipped: {skipped_count}.{log_suffix}",
            )

        return True, f"Batch conversion successful! Converted {converted_count} files."
    except Exception as exc:
        return False, str(exc)
